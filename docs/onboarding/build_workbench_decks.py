#!/usr/bin/env python3
"""Generate two Workbench onboarding decks (PD + Technical) sharing one step spine.

Grounded in source read this session:
- Shell: pages/workbench/[requestId].js (query-string ?tab=&sub= routing; RequireAppAccess 'reviewers')
- Live tabs: Overview, Proposal, Reviewers, Status, Awardee (5 lifecycle tabs are placeholders)
- Reviewers sub-tabs: Find, Invite Reviewers, Track Reviewers (collapsed S280 from 5: Find, Candidates, Invite, Track, Completed)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x1F, 0x29, 0x37)
ACCENT = RGBColor(0x0B, 0x5C, 0x8A)
GRAY = RGBColor(0x55, 0x5F, 0x6B)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

OUT_DIR = os.path.dirname(os.path.abspath(__file__))  # writes the .pptx next to this script


def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def title_slide(prs, title, subtitle, tag):
    s = _blank(prs)
    # accent band
    band = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(7.5))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.6)).text_frame
    tb.word_wrap = True
    p = tb.paragraphs[0]; p.text = title
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = WHITE
    p2 = tb.add_paragraph(); p2.text = subtitle
    p2.font.size = Pt(20); p2.font.color.rgb = RGBColor(0xC9, 0xD4, 0xE0)
    p3 = tb.add_paragraph(); p3.text = tag
    p3.font.size = Pt(13); p3.font.color.rgb = ACCENT; p3.font.bold = True
    return s


def content_slide(prs, kicker, title, bullets, notes=None):
    s = _blank(prs)
    # top accent bar
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.25))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    kt = s.shapes.add_textbox(Inches(0.6), Inches(0.18), Inches(12), Inches(0.4)).text_frame
    kp = kt.paragraphs[0]; kp.text = kicker
    kp.font.size = Pt(12); kp.font.bold = True; kp.font.color.rgb = RGBColor(0x8F, 0xB6, 0xD6)
    tt = s.shapes.add_textbox(Inches(0.6), Inches(0.52), Inches(12.1), Inches(0.7)).text_frame
    tt.word_wrap = True
    tp = tt.paragraphs[0]; tp.text = title
    tp.font.size = Pt(26); tp.font.bold = True; tp.font.color.rgb = WHITE

    body = s.shapes.add_textbox(Inches(0.7), Inches(1.55), Inches(12.0), Inches(5.6)).text_frame
    body.word_wrap = True
    first = True
    for item in bullets:
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item
        note = not isinstance(level, int)   # ("", text) => italic aside
        if note:
            level = 1
        p = body.paragraphs[0] if first else body.add_paragraph()
        first = False
        p.text = text
        p.level = level
        if note:
            p.font.size = Pt(13); p.font.italic = True; p.font.color.rgb = ACCENT
            p.space_before = Pt(8)
        elif level == 0:
            p.font.size = Pt(17); p.font.bold = True; p.font.color.rgb = NAVY
            p.space_before = Pt(10)
        elif level == 1:
            p.font.size = Pt(14); p.font.color.rgb = GRAY; p.space_before = Pt(2)
        else:
            p.font.size = Pt(12); p.font.color.rgb = GRAY
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def section_slide(prs, num, title):
    s = _blank(prs)
    band = s.shapes.add_shape(1, 0, Inches(2.6), prs.slide_width, Inches(2.3))
    band.fill.solid(); band.fill.fore_color.rgb = LIGHT; band.line.fill.background()
    nb = s.shapes.add_textbox(Inches(0.9), Inches(2.75), Inches(2), Inches(1)).text_frame
    np_ = nb.paragraphs[0]; np_.text = num
    np_.font.size = Pt(54); np_.font.bold = True; np_.font.color.rgb = ACCENT
    tb = s.shapes.add_textbox(Inches(3.0), Inches(2.95), Inches(9.4), Inches(1.6)).text_frame
    tb.word_wrap = True
    tp = tb.paragraphs[0]; tp.text = title
    tp.font.size = Pt(30); tp.font.bold = True; tp.font.color.rgb = NAVY
    return s


# ---------------------------------------------------------------------------
# SHARED SPINE — step titles both decks walk in the same order
# ---------------------------------------------------------------------------
SPINE = [
    "What the Workbench is",
    "Getting in & opening a request",
    "The tab strip",
    "Overview tab — the command center",
    "Proposal tab",
    "Reviewers · Find",
    "Reviewers · Candidates (and inviting)",
    "Campaign settings",
    "Sending the invitation",
    "What the reviewer sees",
    "Reviewers · Invite & Track",
    "Quota & winding down",
    "Reviewers · Completed",
    "Status tab",
    "Awardee tab — grantee deliverables",
]

# ===========================================================================
# DECK 1 — PD EDITION
# ===========================================================================
def build_pd():
    prs = new_deck()
    title_slide(
        prs,
        "The Request Workbench",
        "A practical guide for Program Directors",
        "ONBOARDING · PD EDITION · DRAFT v1",
    )
    content_slide(prs, "AGENDA", "What you'll be able to do", [
        "Open any request and see, in one place, where it stands and what to do next.",
        "Run the whole reviewer cycle — find, invite, chase, track — without leaving the request.",
        "Read the proposal, its Phase I documents, and the AI summaries side by side.",
        "Send grantee-deliverable invitations to awardees.",
        ("", "We'll walk one proposal from open → reviewers engaged → reviews in. Same step order Connor's technical deck uses."),
    ])

    section_slide(prs, "1", SPINE[0])
    content_slide(prs, "STEP 1", "What the Workbench is", [
        "One command center per request — not a pile of separate apps.",
        ("", "Before: reviewer-finder, review-manager, and lookups were separate tools you switched between."),
        ("", "Now: open the request once; everything for it lives behind a tab strip."),
        "Request-in-context: every action already knows which proposal, cycle, and program you're on.",
        "You recommend; the board decides. The Workbench never changes the official decision.",
    ])

    section_slide(prs, "2", SPINE[1])
    content_slide(prs, "STEP 2", "Getting in & opening a request", [
        "Sign in with your Microsoft account; access is granted per person.",
        "Open a request and you land on the Overview tab.",
        "The header always shows the request number, title, cycle, program, and institution.",
        ("", "Tip: the page remembers the tab in the URL, so you can bookmark or share a deep link to an exact tab."),
    ])

    section_slide(prs, "3", SPINE[2])
    content_slide(prs, "STEP 3", "The tab strip", [
        "Five tabs are live today:",
        (1, "Overview · Proposal · Reviewers · Status · Awardee"),
        "Five more are placeholders for the rest of the request lifecycle:",
        (1, "Initial Writeup · Reviews · Pre/Site Visit · Final Writeup — they say \"coming in a later update.\""),
        "You'll spend most of your time on the Reviewers tab.",
    ])

    section_slide(prs, "4", SPINE[3])
    content_slide(prs, "STEP 4", "Overview tab — the command center", [
        "A fast answer to \"where is this request and what do I do next?\"",
        "Reviewer funnel at a glance: how many candidates, invited, accepted, materials sent, reviews in.",
        "A \"what next\" hint based on that funnel.",
        "Key request facts: meeting date, amounts, status.",
        ("", "Use it as your daily starting point for each request."),
    ])

    section_slide(prs, "5", SPINE[4])
    content_slide(prs, "STEP 5", "Proposal tab", [
        "Everything about the proposal itself, stacked top to bottom:",
        (1, "Proposal info — PI, Co-PIs, abstract, amounts."),
        (1, "Phase I documents — view/download the submitted files."),
        (1, "AI content — fit rationale, summary, extracted data (Field Primer expands here over time)."),
        ("", "This is your read-the-proposal surface while you decide who should review it."),
    ])

    section_slide(prs, "6", SPINE[5])
    content_slide(prs, "STEP 6", "Reviewers · Find", [
        "Build a pool of potential reviewers for this proposal.",
        "Pull in applicant-recommended reviewers automatically.",
        "Search for additional candidates by expertise; the tool enriches contact details.",
        "Nothing is saved until you choose \"Save as candidates\" — search results are scratch until then.",
    ])

    section_slide(prs, "7", SPINE[6])
    content_slide(prs, "STEP 7", "Reviewers · Candidates (and inviting)", [
        "Your saved shortlist for this request, with the detail you need to decide:",
        (1, "Why they fit, h-index/citations, expertise keywords, Scholar/ORCID/website links."),
        "Select candidates and click Send invitation to invite the ones not yet invited.",
        "Heads-up — there is no longer a \"Re-invite\" button:",
        (1, "Invited people who haven't replied are nudged automatically by a reminder (Step 8/11)."),
    ])

    section_slide(prs, "8", SPINE[7])
    content_slide(prs, "STEP 8", "Campaign settings", [
        "Set the timing once for the whole request; invitations inherit it.",
        ("", "Days to respond — how long a reviewer has to accept (an offset, so later waves get the full window)."),
        ("", "Review due date — the fixed deadline reviews are due."),
        ("", "Reminders — turn the respond-by and review-due nudges on/off and set how many days before."),
        ("", "Desired count — how many accepted reviewers you're aiming for (drives the quota notice)."),
    ])

    section_slide(prs, "9", SPINE[8])
    content_slide(prs, "STEP 9", "Sending the invitation", [
        "The invitation modal pre-fills timing from Campaign settings.",
        "Review the email, then send — each reviewer gets a secure personal accept/decline link.",
        "Low-confidence or missing email addresses are flagged before you send.",
        ("", "Only the invitation goes pre-acceptance — proposal materials are never attached until a reviewer has accepted."),
    ])

    section_slide(prs, "10", SPINE[9])
    content_slide(prs, "STEP 10", "What the reviewer sees", [
        "Reviewers click their link and land on a simple accept/decline page.",
        "Accepting captures the COI & AI-use acknowledgements plus mailing/contact details for honorarium.",
        "They then \"sit tight\" until you release the proposal materials.",
        ("", "Knowing this page helps you answer reviewer questions quickly."),
    ])

    section_slide(prs, "11", SPINE[10])
    content_slide(prs, "STEP 11", "Reviewers · Invite & Track", [
        "Invite sub-tab — accepted reviewers awaiting materials; use Release to reviewers to send the proposal.",
        "Track sub-tab — reviewers in flight, from materials sent through review received.",
        "The system chases for you:",
        (1, "Respond-by reminder nudges invitees who haven't replied."),
        (1, "Review-due reminder nudges accepted reviewers who haven't submitted."),
    ])

    section_slide(prs, "12", SPINE[11])
    content_slide(prs, "STEP 12", "Quota & winding down", [
        "When accepted reviewers reach your desired count, you get a notification.",
        "You decide who (if anyone) to release — a slow-but-wanted senior reviewer is never auto-dropped.",
        "Select still-pending invitees and choose Release as no longer needed:",
        (1, "Sends a polite thank-you and closes their invitation; their reminders stop."),
    ])

    section_slide(prs, "13", SPINE[12])
    content_slide(prs, "STEP 13", "Reviewers · Completed", [
        "Reviewers whose review is in.",
        "Your end-state view of the panel for this request.",
    ])

    section_slide(prs, "14", SPINE[13])
    content_slide(prs, "STEP 14", "Status tab", [
        "A read-only reflection of the proposal's official Dynamics status.",
        "Staff recommend; the board decides; the decision is recorded elsewhere in Dynamics.",
        "The Workbench never edits status here — it's a display surface only.",
    ])

    section_slide(prs, "15", SPINE[14])
    content_slide(prs, "STEP 15", "Awardee tab — grantee deliverables", [
        "For research awards: orchestrate the grantee-deliverables outreach.",
        "Generate the style-guide abstract, confirm recipients (PI in To, liaison in Cc), preview, and send the invite.",
        "Copy website HTML for a single award, or open the combined cycle export.",
    ])

    content_slide(prs, "WRAP-UP", "FAQ & where to get help", [
        "\"Where did the Re-invite button go?\" — replaced by the automatic respond-by reminder.",
        "\"A change isn't showing\" — hard-refresh; the app updates as we ship.",
        "\"A tab says coming soon\" — that lifecycle stage isn't built yet; it's on the roadmap.",
        "Who to ask: your team lead for workflow; Connor for anything technical.",
        ("", "This deck is a draft — tell us what's confusing and we'll tighten it."),
    ])
    return prs


# ===========================================================================
# DECK 2 — TECHNICAL EDITION (Connor)
# ===========================================================================
def build_tech():
    prs = new_deck()
    title_slide(
        prs,
        "The Request Workbench — Under the Hood",
        "Architecture & data walkthrough on the same step spine",
        "ONBOARDING · TECHNICAL EDITION · DRAFT v1",
    )
    content_slide(prs, "ORIENTATION", "How this deck maps to the PD deck", [
        "Same 15-step spine, same order — each step adds the implementation layer beneath it.",
        "Conventions in this deck:",
        (1, "file paths in pages/ shared/ lib/ ; API routes as /api/... ; Dataverse logical names as wmkf_* / akoya_*."),
        "Source of truth, in priority order: live source & probes > Atlas pages > this deck.",
        ("", "Citations were read from source this session; verify against Atlas before quoting in anything durable."),
    ])

    content_slide(prs, "ARCHITECTURE", "The shell", [
        "Per-request shell: pages/workbench/[requestId].js (tier-3 page).",
        "Tab + sub-tab selection is query-string driven: ?tab=reviewers&sub=invite (deep-linkable).",
        "Live tab components:",
        (1, "OverviewTab · ProposalTab · ReviewersTab · StatusTab · AwardeeTab (shared/components/...)."),
        (1, "The other 5 tab keys fall through to a \"coming in a later update\" card."),
        "Consolidation: the Workbench merges the old reviewer-finder + review-manager apps; one app-access grant ('reviewers') now covers them.",
    ])
    content_slide(prs, "ARCHITECTURE", "Access & auth", [
        "Page is wrapped in RequireAppAccess appKey=\"reviewers\".",
        "Per-user app grants live in Dataverse (admin-managed); the same filter gates tab visibility.",
        "Identity comes from the authenticated session — never from request input.",
        "canManage is computed (reviewer-modes) to gate write actions vs. read-only viewers.",
    ])

    section_slide(prs, "1-3", "Shell, entry & tab strip")
    content_slide(prs, "STEPS 1-3", "Open → context → tabs (the data)", [
        "Request context is resolved once by the shell via /api/workbench/resolve-request.",
        (1, "Returns requestId, requestNumber, title, cycleLabel, grantProgram, institution, PD."),
        "context is passed down to every tab component — tabs don't re-fetch identity.",
        "Cycle label derives from the meeting date (J{YY}/D{YY}); used by Awardee export.",
    ])

    section_slide(prs, "4", SPINE[3])
    content_slide(prs, "STEP 4", "Overview tab", [
        "Component: shared/components/workbench/OverviewTab.js (v1, command center).",
        "Data: /api/workbench/reviewer-rollup — counts-only funnel (no person/researcher fan-out).",
        "Server derives the funnel counts + a lightweight \"what next\" hint.",
        ("", "Deferred (no live source yet): writeup status, returned-review state, full isActionableForPD rule set."),
    ])

    section_slide(prs, "5", SPINE[4])
    content_slide(prs, "STEP 5", "Proposal tab", [
        "Component: ProposalTab.js — three stacked sections.",
        (1, "Top + AI render from the already-loaded context (no extra call)."),
        (1, "Documents fetched from /api/workbench/proposal-documents (request-scoped)."),
        "Phase I files download through the scoped blob proxy (no raw blob token to the client).",
        "Field Primer envelope parsed client-side (utils/field-primer-envelope, field-primer-display).",
    ])

    section_slide(prs, "6", SPINE[5])
    content_slide(prs, "STEP 6", "Reviewers · Find", [
        "Component: ReviewerFindPanel (Phase 3) — applicant ingestion + in-panel search.",
        "Search/enrich routes under /api/reviewer-finder/* (discover, enrich-contacts, analyze).",
        "External enrichment is affiliation-disambiguated (serp-contact-service institutionConflicts).",
        "Results are EPHEMERAL until save-candidates writes wmkf_appreviewersuggestion.",
    ])

    section_slide(prs, "7", SPINE[6])
    content_slide(prs, "STEP 7", "Reviewers · Candidates (+ invite)", [
        "Component: ReviewerInvitePanel.js; roster from /api/reviewer-finder/my-candidates?requestId=.",
        "Candidate detail maps off the person row wmkf_potentialreviewers (bibliometrics folded on; sidecar dropped S213).",
        "Invite → InviteEmailModal → render-emails / send-emails with templateType:'invitation'.",
        "S277: the manual \"Re-invite already-invited\" button was REMOVED;",
        (1, "server allowResend re-mint contract retained (shouldSkipDuplicateInvitation) for programmatic re-mint only."),
    ])

    section_slide(prs, "8", SPINE[7])
    content_slide(prs, "STEP 8", "Campaign settings", [
        "Route: /api/review-manager/campaign-config (GET/POST).",
        "Persisted per-request on akoya_request (wave 7-reviewer-engagement, prod 2026-06-21):",
        (1, "wmkf_respondoffsetdays, wmkf_reviewduedate, wmkf_respondreminderenabled/leaddays,"),
        (1, "wmkf_reviewduereminderenabled/leaddays, wmkf_desiredcount, wmkf_quotanotifiedat."),
        "Discrete columns (not a JSON blob) so the cron/sweep can $filter server-side.",
    ])

    section_slide(prs, "9", SPINE[8])
    content_slide(prs, "STEP 9", "Sending the invitation", [
        "render-emails mints the accept/decline magic link ({{externalLink}}); send-emails delivers via Dynamics.",
        "Duplicate-send guard: shouldSkipDuplicateInvitation (lib/utils/reviewer-invite.js).",
        "Attachment safety is SERVER-authoritative:",
        (1, "recipientMayReceiveAttachments gates on recipient wmkf_accepted, NOT caller templateType."),
        "Sets wmkf_invited + wmkf_emailsentat; no reviewstatus bump (pre-acceptance).",
    ])

    section_slide(prs, "10", SPINE[9])
    content_slide(prs, "STEP 10", "What the reviewer sees (portal)", [
        "Page: pages/external/review/[token].js; routes /api/external/review/[token]/{context,respond,upload}.",
        "Token: signed JWT, only its hash stored; re-mints on every email with the link (latest-link-wins).",
        "verify-suggestion-token checks signature/expiry + stored hash + revoked + stored expiry.",
        "Accept (respond.js): COI + AI acks (400 if missing), address/phone (422 unless opt-out); honorarium capture-only this cycle.",
    ])

    section_slide(prs, "11", SPINE[10])
    content_slide(prs, "STEP 11", "Invite & Track + token TTL + reminders", [
        "Release/materials send mints a long-lived token (~review-due + 90d); only accepted reviewers receive it.",
        "Non-responder/invite links cap at review-due + grace (lib/external/reviewer-token-ttl.js via render-emails).",
        "Reminders cron: /api/cron/reviewer-reminders (daily) → lib/services/reviewer-reminder-sweep.js.",
        (1, "respond-by: deadline = emailSentAt + respondOffsetDays - lead; fire-once wmkf_respondremindersentat."),
        (1, "review-due: deadline = reviewDueDate - lead; fire-once via existing wmkf_remindersentat. Claim-before-send (If-Match)."),
    ])

    section_slide(prs, "12", SPINE[11])
    content_slide(prs, "STEP 12", "Quota & winding down", [
        "lib/services/reviewer-quota.js — maybeNotifyQuotaReached called from respond.js AFTER the accept commits.",
        "Count = wmkf_accepted eq true; notify is a conditional null->set of wmkf_quotanotifiedat via If-Match (no double-notify).",
        "Release-as-no-longer-needed: POST /api/review-manager/withdraw-sufficient.",
        (1, "writes withdrawn_sufficient + wmkf_withdrawnsufficientat on still-pending rows only; clears the respond marker."),
    ])

    section_slide(prs, "13-14", "Completed & Status")
    content_slide(prs, "STEPS 13-14", "Completed list & Status reflection", [
        "Completed: reviews returned (review-manager reviewers feed; wmkf_accepted lifecycle).",
        "StatusTab.js: read-only reflection of akoya_requeststatus.",
        (1, "statusClass derived via the value->class map in lib/services/dataverse-export/constants.js."),
        (1, "Unknown status -> UNCLASSIFIED, shown raw (never coerced). Workbench never writes status."),
    ])

    section_slide(prs, "15", SPINE[14])
    content_slide(prs, "STEP 15", "Awardee tab — grantee deliverables", [
        "Component: AwardeeTab.js (chunk 3d); research-grants only.",
        "Endpoints under /api/workbench/grantee-deliverables/*:",
        (1, "generate (abstract), recipients, send-invite, website-html, cycle-export."),
        "send-invite mints a stateless grantee magic link server-side; PI in To, liaison in Cc; flips Drafted->Invited once.",
    ])

    content_slide(prs, "APPENDIX", "Data & schema", [
        "Reviewer engagement wave: lib/dataverse/schema/wave7-reviewer-engagement/ (9 fields, prod 2026-06-21).",
        "Per-request config → akoya_request; per-reviewer markers → wmkf_appreviewersuggestion.",
        "Candidate/person identity → wmkf_potentialreviewers (bibliometrics on the person).",
        "Status taxonomy → lib/services/dataverse-export/constants.js.",
        ("", "New columns carry NO Power Automate trigger (verified post-deploy)."),
    ])
    content_slide(prs, "APPENDIX", "Automation & ops", [
        "Crons in pages/api/cron/: reviewer-reminders (daily), sweep-stale-invites (status bookkeeping at meeting date).",
        "Token cap (access gate at review-due) and the stale-invite sweep (status at meeting date) are DIFFERENT gates — both stay.",
        "Cron auth: Vercel CRON_SECRET on all /api/cron/* routes.",
        "Reminder/quota writes use If-Match claim-before-send for at-most-once semantics.",
    ])
    content_slide(prs, "POINTERS", "Source of truth (don't let this deck rot)", [
        "Spec: docs/REVIEWER_ENGAGEMENT_SPEC.md (status: IMPLEMENTED, all 4 phases LIVE S275).",
        "Atlas: docs/atlas/dataverse-akoya-request.md, dataverse-wmkf-appreviewersuggestion.md.",
        "Agent wiki: docs/agent-wiki/topics/reviewer-workbench-lifecycle.md (+ reviewer-* topics).",
        "Build plans: REQUEST_WORKBENCH_BUILD_PLAN.md, WORKBENCH_PROPOSAL_TAB_BUILD_PLAN.md.",
        ("", "When behavior changes, update the spec/Atlas/wiki first; regenerate this deck from them."),
    ])
    return prs


def main():
    os.makedirs(OUT_DIR, exist_ok=True)
    build_pd().save(os.path.join(OUT_DIR, "Workbench_Onboarding_PD.pptx"))
    build_tech().save(os.path.join(OUT_DIR, "Workbench_Onboarding_Technical.pptx"))
    for f in sorted(os.listdir(OUT_DIR)):
        print("wrote", os.path.join(OUT_DIR, f))


if __name__ == "__main__":
    main()
